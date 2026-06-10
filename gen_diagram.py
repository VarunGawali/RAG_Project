"""
gen_diagram.py
Generates architecture_diagram.pptx for Contract360 system.
Two slides: Slide 1 = Ingestion Path, Slide 2 = Query Path.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── helpers ──────────────────────────────────────────────────────────────────

def hex2rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

BG       = hex2rgb("#0D1117")
WHITE    = hex2rgb("#FFFFFF")
ARROW_C  = hex2rgb("#8B949E")
LABEL_C  = hex2rgb("#C9D1D9")

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, x, y, w, h, fill_hex, title_text, body_lines=None,
            title_size=11, body_size=8.5, radius=0.12):
    """Add a rounded-rectangle box with a title line and optional body lines."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    # set corner rounding (0..1 range, ~0.05 looks subtle)
    try:
        shape.adjustments[0] = 0.05
    except Exception:
        pass
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = hex2rgb(fill_hex)

    line = shape.line
    line.color.rgb = hex2rgb(fill_hex)
    line.width = Pt(0.5)

    tf = shape.text_frame
    tf.word_wrap = True

    # Title paragraph
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    run0 = p0.add_run()
    run0.text = title_text
    run0.font.bold = True
    run0.font.size = Pt(title_size)
    run0.font.color.rgb = WHITE
    run0.font.name = "Calibri"

    # Body paragraphs
    if body_lines:
        for line_text in body_lines:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line_text
            run.font.size = Pt(body_size)
            run.font.color.rgb = WHITE
            run.font.name = "Calibri"

    # inner margin
    tf.margin_left   = Inches(0.07)
    tf.margin_right  = Inches(0.07)
    tf.margin_top    = Inches(0.04)
    tf.margin_bottom = Inches(0.04)

    return shape


def add_arrow(slide, x1, y1, x2, y2, label=None, label_size=8):
    """Draw a vertical or horizontal connector arrow."""
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    # Use a line connector
    cx = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    cx.line.color.rgb = ARROW_C
    cx.line.width = Pt(1.5)

    # arrowhead at end
    ln = cx.line._ln
    tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", "none")
    headEnd = etree.SubElement(ln, qn("a:headEnd"))
    headEnd.set("type", "arrow")
    headEnd.set("w", "med")
    headEnd.set("len", "med")

    if label:
        # place a small text box near midpoint
        mx = (x1 + x2) / 2
        my = (y1 + y2) / 2
        tb = slide.shapes.add_textbox(
            Inches(mx + 0.05), Inches(my - 0.12), Inches(1.0), Inches(0.22)
        )
        tf = tb.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        run = p.runs[0]
        run.font.size = Pt(label_size)
        run.font.color.rgb = LABEL_C
        run.font.name = "Calibri"
        run.font.italic = True

    return cx


def add_label(slide, x, y, w, h, text, size=10, bold=False, color=None,
              align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color if color else WHITE
    run.font.name = "Calibri"
    return tb


def add_divider_label(slide, y, text):
    """Full-width section label."""
    add_label(slide, 0.2, y, 12.9, 0.28, text, size=10, bold=True,
              color=hex2rgb("#E6EDF3"))


# ── colour constants ──────────────────────────────────────────────────────────
C_BLUE      = "#1F6FEB"
C_LBLUE     = "#388BFD"
C_PURPLE    = "#8957E5"
C_GREEN     = "#238636"
C_RED       = "#F78166"

# ── slide dimensions ──────────────────────────────────────────────────────────
W = 13.33   # inches
H = 7.5

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)

blank_layout = prs.slide_layouts[6]  # blank


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — INGESTION PATH
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, BG)

# Slide title
add_label(s1, 0.2, 0.08, 12.9, 0.30,
          "Contract360 — System Architecture  |  Slide 1: Ingestion Path",
          size=11, bold=True, color=hex2rgb("#E6EDF3"))

# ── 1. React UI ───────────────────────────────────────────────────────────────
RUI_X, RUI_Y, RUI_W, RUI_H = 0.25, 0.45, 12.83, 0.72
add_box(s1, RUI_X, RUI_Y, RUI_W, RUI_H, C_BLUE,
        "React (Vite) UI",
        ["chat interface  •  contract selector  •  upload panel",
         "SSE token streaming  •  citation cards  •  sidebar",
         "4 s polling for ingestion job status"])

# arrow down
add_arrow(s1, W/2, RUI_Y + RUI_H, W/2, 1.38, label="REST + SSE")

# ── 2. FastAPI ────────────────────────────────────────────────────────────────
API_X, API_Y, API_W, API_H = 0.25, 1.38, 12.83, 0.80
add_box(s1, API_X, API_Y, API_W, API_H, C_BLUE,
        "FastAPI API",
        ["POST /ingest  •  GET /contracts  •  DELETE /contracts/{id}",
         "POST /sessions  •  GET /sessions/{id}/history",
         "POST /sessions/{id}/ask/stream  •  GET /health"])

# two arrows from FastAPI
SPLIT_Y = API_Y + API_H
LEFT_X  = 3.0
RIGHT_X = 10.0

add_arrow(s1, LEFT_X,  SPLIT_Y, LEFT_X,  2.48)
add_arrow(s1, RIGHT_X, SPLIT_Y, RIGHT_X, 2.48)

# ── 3a. Ingestion Worker (left) ───────────────────────────────────────────────
IW_X, IW_Y, IW_W, IW_H = 0.25, 2.48, 6.0, 0.62
add_box(s1, IW_X, IW_Y, IW_W, IW_H, C_LBLUE,
        "Ingestion Worker",
        ["Downloads PDF from Blob  •  invokes pipeline stages",
         "writes job status → Cosmos"])

# ── 3b. Query + Session Service (right) ──────────────────────────────────────
QS_X, QS_Y, QS_W, QS_H = 7.08, 2.48, 6.0, 0.62
add_box(s1, QS_X, QS_Y, QS_W, QS_H, C_LBLUE,
        "Query + Session Service",
        ["routes /ask/stream  •  loads history  •  calls retrieval pipeline"])

# arrow from Q+S down to Cosmos DB
COSMOS_Y = 3.32
add_arrow(s1, RIGHT_X, QS_Y + QS_H, RIGHT_X, COSMOS_Y)

# ── Cosmos DB (right) ─────────────────────────────────────────────────────────
CD_X, CD_Y, CD_W, CD_H = 7.08, COSMOS_Y, 6.0, 0.60
add_box(s1, CD_X, CD_Y, CD_W, CD_H, C_RED,
        "Cosmos DB (NoSQL)",
        ["sessions  •  messages  •  ingestion jobs"])

# arrows down to Ingestion Pipeline
PIPE_Y = 4.12
add_arrow(s1, LEFT_X,  IW_Y + IW_H, LEFT_X,  PIPE_Y)
add_arrow(s1, RIGHT_X, CD_Y + CD_H, RIGHT_X, PIPE_Y)

# ── Ingestion Pipeline ────────────────────────────────────────────────────────
IP_X, IP_Y, IP_W, IP_H = 0.25, PIPE_Y, 12.83, 1.28
add_box(s1, IP_X, IP_Y, IP_W, IP_H, C_GREEN,
        "Ingestion Pipeline",
        ["1. Download PDF from Blob    2. Parse → Azure Document Intelligence    3. Build clause tree",
         "4. Chunk text    5. Batch embed (Azure OpenAI)    6. Generate summary (GPT-4o)",
         "7. Persist artifacts to Blob    8. Upload to Azure AI Search",
         "9. Extract KG entities (GPT-4o)    10. Resolve entities (normalize → dedupe → canonicalize)",
         "11. Write 2-tier graph to Cosmos Gremlin"])

# three arrows down to storage boxes
STORE_Y = PIPE_Y + IP_H + 0.10
SB_W = 3.8
GAP  = 0.62
S1_X = 0.52
S2_X = S1_X + SB_W + GAP / 2 - 0.15
S3_X = S2_X + SB_W + GAP / 2 - 0.15

add_arrow(s1, S1_X + SB_W/2, PIPE_Y + IP_H, S1_X + SB_W/2, STORE_Y)
add_arrow(s1, S2_X + SB_W/2, PIPE_Y + IP_H, S2_X + SB_W/2, STORE_Y)
add_arrow(s1, S3_X + SB_W/2, PIPE_Y + IP_H, S3_X + SB_W/2, STORE_Y)

SB_H = 0.78
add_box(s1, S1_X, STORE_Y, SB_W, SB_H, C_RED,
        "Azure Blob Storage",
        ["PDFs  •  tree.json", "extraction artifacts"])

add_box(s1, S2_X, STORE_Y, SB_W, SB_H, C_RED,
        "Azure AI Search",
        ["BM25 + vector  •  hybrid index", "clause chunks + embeddings"])

add_box(s1, S3_X, STORE_Y, SB_W, SB_H, C_RED,
        "Cosmos DB Gremlin",
        ["2-tier KG  •  Mention nodes", "CanonicalEntity  •  RESOLVED_AS"])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — QUERY PATH
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, BG)

add_label(s2, 0.2, 0.08, 12.9, 0.30,
          "Contract360 — System Architecture  |  Slide 2: Query Path",
          size=11, bold=True, color=hex2rgb("#E6EDF3"))

# ── Query Understanding ───────────────────────────────────────────────────────
QU_X, QU_Y, QU_W, QU_H = 0.25, 0.50, 12.83, 1.05
add_box(s2, QU_X, QU_Y, QU_W, QU_H, C_PURPLE,
        "Query Understanding",
        ["1. Load chat history from Cosmos (last 20 msgs)",
         "2. Summary shortcut — skip retrieval if summarize intent detected",
         "3. Scope resolution — narrow to contracts named in question",
         "4. Route: LLM classifier  →  tree / graph / hybrid  (comparison → hybrid)"])

# two arrows
ROUTE_Y = QU_Y + QU_H + 0.08
TR_X = 3.0
GR_X = 10.33

add_arrow(s2, TR_X, QU_Y + QU_H, TR_X, ROUTE_Y)
add_arrow(s2, GR_X, QU_Y + QU_H, GR_X, ROUTE_Y)

# ── Tree Route ────────────────────────────────────────────────────────────────
TRB_X, TRB_Y, TRB_W, TRB_H = 0.25, ROUTE_Y, 6.0, 0.95
add_box(s2, TRB_X, TRB_Y, TRB_W, TRB_H, C_LBLUE,
        "Tree Route  —  SemanticRetriever",
        ["keyword extract",
         "→ parent / child chunk expansion",
         "→ Azure Search BM25 + vector + rerank"])

# ── Graph Route ───────────────────────────────────────────────────────────────
GRB_X, GRB_Y, GRB_W, GRB_H = 7.08, ROUTE_Y, 6.0, 0.95
add_box(s2, GRB_X, GRB_Y, GRB_W, GRB_H, C_LBLUE,
        "Graph Route",
        ["Phase 1: entity link → RESOLVED_AS → obligations (Gremlin)",
         "Phase 2 fallback: vector search → clause IDs → graph"])

# merge arrows to Comparison Branch
COMP_Y = ROUTE_Y + TRB_H + 0.10
add_arrow(s2, TR_X, TRB_Y + TRB_H, TR_X, COMP_Y)
add_arrow(s2, GR_X, GRB_Y + GRB_H, GR_X, COMP_Y)

# horizontal merge line
add_arrow(s2, TR_X, COMP_Y, GR_X, COMP_Y)

MID_X = (TR_X + GR_X) / 2

# ── Comparison Branch ─────────────────────────────────────────────────────────
CB_X, CB_Y, CB_W, CB_H = 2.5, COMP_Y, 8.33, 0.80
add_box(s2, CB_X, CB_Y, CB_W, CB_H, C_PURPLE,
        "Comparison Branch  (2+ contracts + comparison intent)",
        ["format_comparison_result()",
         "→ CONTRACT A / CONTRACT B side-by-side table"])

# arrow to Answer Generation
AG_Y = CB_Y + CB_H + 0.10
add_arrow(s2, MID_X, CB_Y + CB_H, MID_X, AG_Y)

# ── Answer Generation ─────────────────────────────────────────────────────────
AG_X, AG_W, AG_H = 0.25, 12.83, 1.45
add_box(s2, AG_X, AG_Y, AG_W, AG_H, C_GREEN,
        "Answer Generation + Grounding",
        ["•  Inject [S#] citations into context",
         "•  GPT-4o generates answer with inline [S#] markers",
         "•  Extract only cited sources → grounded citations list",
         "•  Generate 3 follow-up suggestions",
         "•  Persist answer + citations + follow-ups to Cosmos",
         "•  Stream response via SSE word-by-word",
         "•  Final SSE event:  { done, citations, follow_ups }"])

# ── Footnote ─────────────────────────────────────────────────────────────────
add_label(s2, 0.2, H - 0.28, 12.9, 0.22,
          "All LLM calls → Azure OpenAI (GPT-4o)  •  Embeddings → text-embedding-3-large  •  Search → Azure AI Search hybrid",
          size=7.5, color=hex2rgb("#8B949E"), align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "/home/user/RAG_Project_EY/architecture_diagram.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
