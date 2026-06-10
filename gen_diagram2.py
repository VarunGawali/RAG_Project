"""
gen_diagram2.py — Contract360 Architecture Diagram (manager-facing, single slide)
Generates architecture_diagram2.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import pptx.oxml as oxml
from lxml import etree
import copy

# ── helpers ──────────────────────────────────────────────────────────────────

def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_rounded_rect(slide, x, y, w, h, fill_hex, title_text, detail_lines=None,
                     title_size=10, detail_size=8, border_hex='#21262D',
                     corner_radius_emu=None):
    """Add a rounded rectangle with title + optional detail lines."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    # corner radius — adjust the adj XML attribute
    if corner_radius_emu is None:
        # default ~0.08" expressed as ratio of shorter dimension * 100000
        shorter = min(w, h)
        ratio = min(0.08 / shorter, 0.5)
        adj_val = int(ratio * 100000)
    else:
        adj_val = corner_radius_emu

    sp = shape.element
    # find or create prstGeom/avLst/gd
    prstGeom = sp.find('.//' + qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        # remove existing gd
        for gd in avLst.findall(qn('a:gd')):
            avLst.remove(gd)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {adj_val}')

    # fill
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    # border
    shape.line.color.rgb = rgb(border_hex)
    shape.line.width = Pt(0.5)

    # text
    tf = shape.text_frame
    tf.word_wrap = True

    # vertical centering
    from pptx.enum.text import PP_ALIGN
    tf.auto_size = None

    # title paragraph
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    run = p0.add_run()
    run.text = title_text
    run.font.bold = True
    run.font.size = Pt(title_size)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = 'Calibri'

    if detail_lines:
        for dline in detail_lines:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = dline
            r.font.bold = False
            r.font.size = Pt(detail_size)
            r.font.color.rgb = RGBColor(0xC9, 0xD1, 0xD9)
            r.font.name = 'Calibri'

    # vertical alignment — middle
    from pptx.oxml.ns import nsmap
    txBody = tf._txBody
    txBody.set('anchor', 'ctr')

    return shape


def add_label(slide, x, y, w, h, text, size=8, bold=False, color_hex='#8B949E',
              align=PP_ALIGN.CENTER):
    """Add a transparent text label."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.bold = bold
    r.font.size = Pt(size)
    r.font.color.rgb = rgb(color_hex)
    r.font.name = 'Calibri'
    return txBox


def add_arrow(slide, x1, y1, x2, y2, color_hex='#8B949E', width_pt=1.5):
    """Draw a straight connector arrow from (x1,y1) to (x2,y2) in inches."""
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = rgb(color_hex)
    connector.line.width = Pt(width_pt)
    # add arrowhead at end
    ln = connector.line._ln
    tailEnd = ln.find(qn('a:tailEnd'))
    if tailEnd is None:
        tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'none')
    headEnd = ln.find(qn('a:headEnd'))
    if headEnd is None:
        headEnd = etree.SubElement(ln, qn('a:headEnd'))
    headEnd.set('type', 'triangle')
    headEnd.set('w', 'med')
    headEnd.set('len', 'med')
    return connector


def add_arrow_label(slide, x, y, text, size=7.5):
    add_label(slide, x, y, 1.2, 0.22, text, size=size, color_hex='#8B949E')


# ── main ─────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb('#0D1117')

    ARROW = '#6E7681'
    BORDER = '#30363D'

    # ── 1. TITLE BAR ──────────────────────────────────────────────────────────
    title_bar = slide.shapes.add_shape(
        __import__('pptx.enum.shapes', fromlist=['MSO_SHAPE']).MSO_SHAPE.RECTANGLE,
        Inches(0.0), Inches(0.0), Inches(13.33), Inches(0.33)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = rgb('#161B22')
    title_bar.line.fill.background()

    tf = title_bar.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'Contract360  —  System Architecture'
    r.font.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.name = 'Calibri'
    tf._txBody.set('anchor', 'ctr')

    # ── ZONE LABELS ───────────────────────────────────────────────────────────
    add_label(slide, 0.2,  0.35, 5.5, 0.22, '◀  INGESTION PATH', size=7.5,
              bold=True, color_hex='#3FB950', align=PP_ALIGN.LEFT)
    add_label(slide, 6.2,  0.35, 6.8, 0.22, 'QUERY PATH  ▶', size=7.5,
              bold=True, color_hex='#79C0FF', align=PP_ALIGN.RIGHT)

    # vertical divider (subtle dashed line)
    from pptx.enum.shapes import MSO_CONNECTOR
    div = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(6.0), Inches(0.35), Inches(6.0), Inches(6.42)
    )
    div.line.color.rgb = rgb('#21262D')
    div.line.width = Pt(1.0)
    # dashed
    ln = div.line._ln
    prstDash = etree.SubElement(ln, qn('a:prstDash'))
    prstDash.set('val', 'sysDash')

    # ── 2. REACT UI (full width) ──────────────────────────────────────────────
    add_rounded_rect(slide, 0.2, 0.4, 12.93, 0.52,
                     '#1F6FEB',
                     'React (Vite) UI',
                     ['Chat  •  Contract Selector  •  Upload Panel  •  SSE Streaming  •  Citation Cards'],
                     title_size=11, detail_size=8.5, border_hex=BORDER)

    # ── arrow: React → FastAPI ────────────────────────────────────────────────
    add_arrow(slide, 6.665, 0.92, 6.665, 1.05, ARROW)
    add_arrow_label(slide, 6.85, 0.9, 'REST + SSE', size=7)

    # ── 3. FASTAPI (full width) ───────────────────────────────────────────────
    add_rounded_rect(slide, 0.2, 1.05, 12.93, 0.52,
                     '#1158C7',
                     'FastAPI API Layer',
                     ['/ingest  •  /contracts  •  /sessions  •  /ask/stream  •  /health  •  JWT auth  •  CORS'],
                     title_size=11, detail_size=8.5, border_hex=BORDER)

    # ── arrows: FastAPI → left (Ingestion Worker) and → right (Query+Session) ─
    # left fork: start from ~x=1.9 (center of ingestion worker box), top of box at y=1.57
    add_arrow(slide, 1.9,  1.57, 1.9,  1.75, ARROW)
    add_arrow(slide, 8.7,  1.57, 8.7,  1.75, ARROW)

    # ── 4. INGESTION WORKER (left) ────────────────────────────────────────────
    add_rounded_rect(slide, 0.2, 1.75, 3.4, 0.52,
                     '#0D419D',
                     'Ingestion Worker',
                     ['worker.py  •  async job processor  •  status tracking'],
                     title_size=10, detail_size=8, border_hex=BORDER)

    # ── 5. QUERY + SESSION SERVICE (right) ────────────────────────────────────
    add_rounded_rect(slide, 6.2, 1.75, 5.0, 0.52,
                     '#0D419D',
                     'Query + Session Service',
                     ['query_service.py  •  session CRUD  •  contract scoping'],
                     title_size=10, detail_size=8, border_hex=BORDER)

    # ── arrow: Query+Session → Cosmos NoSQL ───────────────────────────────────
    add_arrow(slide, 8.7, 2.27, 8.7, 2.45, ARROW)

    # ── 6. COSMOS DB NoSQL (right) ────────────────────────────────────────────
    add_rounded_rect(slide, 6.2, 2.45, 5.0, 0.58,
                     '#B45309',
                     'Cosmos DB  (NoSQL)',
                     ['Chat sessions  •  Messages  •  Ingestion jobs  •  Contract metadata'],
                     title_size=10, detail_size=8, border_hex=BORDER)

    # ── arrows: Worker → Ingestion Pipeline; Cosmos → Query Router ────────────
    add_arrow(slide, 1.9,  2.27, 1.9,  3.2,  ARROW)   # Worker down
    add_arrow(slide, 8.7,  3.03, 8.7,  3.2,  ARROW)   # Cosmos → Query Router

    # ── 7. INGESTION PIPELINE (left+center) ───────────────────────────────────
    add_rounded_rect(slide, 0.2, 3.2, 5.55, 0.95,
                     '#238636',
                     'Ingestion Pipeline',
                     [
                         'PDF Parse (Doc Intelligence)  →  Clause Tree  →  Chunk  →  Batch Embed (text-embedding-3-large)',
                         '→  Summary (GPT-4o)  →  Search Index  →  KG Extract  →  Entity Resolve  →  Graph Write'
                     ],
                     title_size=10, detail_size=7.5, border_hex=BORDER)

    # ── 8. QUERY ROUTER (right) ───────────────────────────────────────────────
    add_rounded_rect(slide, 6.2, 3.2, 6.93, 0.62,
                     '#6E40C9',
                     'Query Router',
                     ['LLM classifier  •  scope resolution  •  summary shortcut  •  hybrid / graph / compare'],
                     title_size=10, detail_size=8, border_hex=BORDER)

    # ── arrows: Query Router → Tree Route and → Graph Route ──────────────────
    add_arrow(slide, 7.5,  3.82, 7.5,  3.97, ARROW)   # → Tree Route
    add_arrow(slide, 11.3, 3.82, 11.3, 3.97, ARROW)   # → Graph Route

    # fork line from Query Router bottom center → two branches
    # horizontal connector at y=3.9 from x=7.5 to x=11.3
    h_fork = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(7.5), Inches(3.9), Inches(11.3), Inches(3.9)
    )
    h_fork.line.color.rgb = rgb(ARROW)
    h_fork.line.width = Pt(1.2)

    # ── 9. TREE ROUTE (right, left sub-column) ────────────────────────────────
    add_rounded_rect(slide, 6.2, 3.97, 3.0, 0.82,
                     '#388BFD',
                     'Tree Route',
                     [
                         'SemanticRetriever  •  Azure Search',
                         'BM25 + vector  •  parent/child expansion'
                     ],
                     title_size=9.5, detail_size=7.5, border_hex=BORDER)

    # ── 10. GRAPH ROUTE (right, right sub-column) ─────────────────────────────
    add_rounded_rect(slide, 9.55, 3.97, 3.58, 0.82,
                     '#388BFD',
                     'Graph Route',
                     [
                         'Entity linking → RESOLVED_AS',
                         'Gremlin traversal  •  vector fallback'
                     ],
                     title_size=9.5, detail_size=7.5, border_hex=BORDER)

    # ── arrows: Tree/Graph → Hybrid ───────────────────────────────────────────
    add_arrow(slide, 7.7,  4.79, 7.7,  4.97, ARROW)
    add_arrow(slide, 11.34, 4.79, 11.34, 4.97, ARROW)
    # horizontal merge line
    h_merge = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(7.7), Inches(4.9), Inches(11.34), Inches(4.9)
    )
    h_merge.line.color.rgb = rgb(ARROW)
    h_merge.line.width = Pt(1.2)

    # ── 11. HYBRID MERGE + COMPARISON ─────────────────────────────────────────
    add_rounded_rect(slide, 6.2, 4.97, 6.93, 0.62,
                     '#6E40C9',
                     'Hybrid Merge  +  Comparison Branch',
                     ['Both routes merged  •  CONTRACT A / B side-by-side when 2+ contracts scoped'],
                     title_size=10, detail_size=8, border_hex=BORDER)

    # ── arrow: Hybrid → Answer Gen ────────────────────────────────────────────
    add_arrow(slide, 9.665, 5.59, 9.665, 5.75, ARROW)

    # ── 12. ANSWER GENERATION ─────────────────────────────────────────────────
    add_rounded_rect(slide, 6.2, 5.75, 6.93, 0.62,
                     '#6E40C9',
                     'Answer Generation',
                     ['GPT-4o  •  [S#] grounding  •  citations  •  SSE streaming back to UI'],
                     title_size=10, detail_size=8, border_hex=BORDER)

    # ── arrows: Ingestion Pipeline → Storage layer ────────────────────────────
    add_arrow(slide, 1.1,  4.15, 1.1,  6.45, ARROW)   # → Blob
    add_arrow(slide, 3.0,  4.15, 3.0,  6.45, ARROW)   # → Search
    add_arrow(slide, 4.8,  4.15, 4.8,  6.45, ARROW)   # → Gremlin

    # ── arrows: Query routes → Storage layer ─────────────────────────────────
    add_arrow(slide, 7.7,  4.79, 3.0,  6.45, ARROW)   # Tree → Search (long)
    add_arrow(slide, 11.34, 4.79, 4.8, 6.45, ARROW)   # Graph → Gremlin (long)

    # ── 13. STORAGE BAR BACKGROUND ────────────────────────────────────────────
    from pptx.enum.shapes import MSO_SHAPE
    storage_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.0), Inches(6.42), Inches(13.33), Inches(1.08)
    )
    storage_bg.fill.solid()
    storage_bg.fill.fore_color.rgb = rgb('#1A1F2E')
    storage_bg.line.fill.background()

    add_label(slide, 0.2, 6.43, 3.0, 0.22, 'AZURE STORAGE & SERVICES LAYER',
              size=7.5, bold=True, color_hex='#E3B341', align=PP_ALIGN.LEFT)

    # 4 storage boxes
    storage_boxes = [
        (0.2,  6.67, 2.8,  'Azure Blob Storage',   'Raw PDFs  •  chunked text  •  processed docs'),
        (3.2,  6.67, 3.3,  'Azure AI Search',       'BM25 + vector index  •  hybrid retrieval'),
        (6.7,  6.67, 3.3,  'Cosmos DB  (Gremlin)',  'Knowledge Graph  •  entities  •  relationships'),
        (10.2, 6.67, 2.93, 'Azure OpenAI',          'GPT-4o (chat)  •  text-embedding-3-large'),
    ]
    for sx, sy, sw, stitle, sdetail in storage_boxes:
        add_rounded_rect(slide, sx, sy, sw, 0.72,
                         '#B45309',
                         stitle, [sdetail],
                         title_size=9, detail_size=7.5, border_hex='#92400E')

    # ── save ──────────────────────────────────────────────────────────────────
    out = '/home/user/RAG_Project_EY/architecture_diagram2.pptx'
    prs.save(out)
    print(f'Saved → {out}')


if __name__ == '__main__':
    build()
